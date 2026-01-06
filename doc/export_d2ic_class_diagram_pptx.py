from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


@dataclass(frozen=True)
class SvgRectStyle:
    fill: str | None
    stroke: str | None
    stroke_width_px: float | None
    stroke_dasharray: str | None


@dataclass(frozen=True)
class SvgBox:
    x: float
    y: float
    w: float
    h: float


@dataclass(frozen=True)
class Node:
    key: str
    box: SvgBox
    lines: list[tuple[str, bool]]  # (text, bold)
    style: SvgRectStyle


@dataclass(frozen=True)
class Cluster:
    key: str
    label: str
    box: SvgBox


@dataclass(frozen=True)
class Edge:
    src: str
    dst: str
    label: str | None
    kind: str  # "solid" | "dotted"


SVG_NS = "http://www.w3.org/2000/svg"
NS = {"svg": SVG_NS}


def _parse_viewbox(root: ET.Element) -> SvgBox:
    vb = root.attrib.get("viewBox", "")
    parts = [p for p in vb.replace(",", " ").split() if p.strip()]
    if len(parts) != 4:
        raise ValueError(f"Unexpected viewBox: {vb!r}")
    x, y, w, h = (float(parts[0]), float(parts[1]), float(parts[2]), float(parts[3]))
    return SvgBox(x=x, y=y, w=w, h=h)


_TRANSLATE_RE = re.compile(r"translate\(\s*([-\d.]+)\s*,\s*([-\d.]+)\s*\)")


def _parse_translate(transform: str | None) -> tuple[float, float]:
    if not transform:
        return 0.0, 0.0
    m = _TRANSLATE_RE.search(transform)
    if not m:
        return 0.0, 0.0
    return float(m.group(1)), float(m.group(2))


def _parse_style(style: str | None) -> SvgRectStyle:
    if not style:
        return SvgRectStyle(fill=None, stroke=None, stroke_width_px=None, stroke_dasharray=None)
    props: dict[str, str] = {}
    for chunk in style.split(";"):
        if ":" not in chunk:
            continue
        k, v = chunk.split(":", 1)
        k = k.strip()
        v = v.strip().replace("!important", "").strip()
        if k:
            props[k] = v
    sw = None
    if "stroke-width" in props:
        try:
            sw = float(props["stroke-width"].replace("px", "").strip())
        except ValueError:
            sw = None
    return SvgRectStyle(
        fill=props.get("fill"),
        stroke=props.get("stroke"),
        stroke_width_px=sw,
        stroke_dasharray=props.get("stroke-dasharray"),
    )


def _hex_to_rgb(hex_color: str) -> RGBColor:
    c = hex_color.strip()
    if c.startswith("#"):
        c = c[1:]
    if len(c) != 6:
        raise ValueError(f"Unsupported color: {hex_color!r}")
    return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))


def _dash_from_svg(dasharray: str | None) -> MSO_LINE_DASH_STYLE | None:
    if not dasharray:
        return None
    d = dasharray.strip()
    if d in ("0", "none"):
        return None
    # Heuristic mapping for this project:
    # - "6 4" => dashed
    # - "2 3" => dotted
    nums = []
    for part in d.replace(",", " ").split():
        try:
            nums.append(float(part))
        except ValueError:
            pass
    if not nums:
        return None
    if len(nums) >= 2 and nums[0] <= 2.5:
        return MSO_LINE_DASH_STYLE.ROUND_DOT
    return MSO_LINE_DASH_STYLE.DASH


def _iter_tspans(text_el: ET.Element) -> Iterable[ET.Element]:
    for tspan in list(text_el):
        if tspan.tag == f"{{{SVG_NS}}}tspan":
            yield tspan


def _extract_lines(text_el: ET.Element) -> list[tuple[str, bool]]:
    lines: list[tuple[str, bool]] = []
    for tspan in _iter_tspans(text_el):
        txt = (tspan.text or "").replace("\u00A0", "").strip("\n")
        bold = (tspan.attrib.get("font-weight", "").lower() == "bold")
        # Keep intentional empty lines as spacer
        if txt.strip() == "":
            lines.append(("", False))
        else:
            lines.append((txt, bold))
    while lines and lines[-1][0] == "":
        lines.pop()
    return lines


def load_nodes(svg_path: Path) -> tuple[SvgBox, list[Cluster], list[Node]]:
    root = ET.fromstring(svg_path.read_text(encoding="utf-8"))
    viewbox = _parse_viewbox(root)

    clusters: list[Cluster] = []
    for c in root.findall('.//svg:g[@class="cluster"]', NS):
        cid = c.attrib.get("id", "").strip()
        rect = c.find("./svg:rect", NS)
        label_el = c.find('.//svg:text', NS)
        if not cid or rect is None or label_el is None:
            continue
        box = SvgBox(
            x=float(rect.attrib.get("x", "0")),
            y=float(rect.attrib.get("y", "0")),
            w=float(rect.attrib.get("width", "0")),
            h=float(rect.attrib.get("height", "0")),
        )
        label = "".join(label_el.itertext()).strip()
        clusters.append(Cluster(key=cid, label=label, box=box))

    nodes: list[Node] = []
    for g in root.findall('.//svg:g[@class="node default"]', NS):
        gid = g.attrib.get("id", "").strip()
        tx, ty = _parse_translate(g.attrib.get("transform"))
        rect = g.find("./svg:rect", NS)
        label_group = g.find('./svg:g[@class="label"]', NS)
        text_el = label_group.find(".//svg:text", NS) if label_group is not None else None
        if rect is None or text_el is None:
            continue
        x = tx + float(rect.attrib.get("x", "0"))
        y = ty + float(rect.attrib.get("y", "0"))
        w = float(rect.attrib.get("width", "0"))
        h = float(rect.attrib.get("height", "0"))
        style = _parse_style(rect.attrib.get("style"))
        lines = _extract_lines(text_el)

        # Heuristic "key" for edges: prefer the first bold line, else first non-empty line, else gid
        key = None
        for t, bold in lines:
            if bold and t.strip():
                key = t.strip()
                break
        if key is None:
            for t, _bold in lines:
                if t.strip():
                    key = t.strip().strip("«»")
                    break
        if key is None:
            key = gid

        nodes.append(Node(key=key, box=SvgBox(x=x, y=y, w=w, h=h), lines=lines, style=style))

    return viewbox, clusters, nodes


_EDGE_RE = re.compile(
    r"^\s*"
    r"(?P<src>[A-Za-z_][A-Za-z0-9_]*)"
    r"\s*(?P<arrow>-->|-\.->)\s*"
    r"(?:(?:\|(?P<label>[^|]+)\|)\s*)?"
    r"(?P<dst>[A-Za-z_][A-Za-z0-9_]*)"
    r"\s*$"
)


def load_edges(mmd_path: Path) -> list[Edge]:
    edges: list[Edge] = []
    for raw in mmd_path.read_text(encoding="utf-8").splitlines():
        line = raw.strip()
        if not line or line.startswith("%%") or line.startswith("flowchart"):
            continue
        m = _EDGE_RE.match(line)
        if not m:
            continue
        src = m.group("src")
        dst = m.group("dst")
        arrow = m.group("arrow")
        label = m.group("label").strip() if m.group("label") else None
        kind = "dotted" if ".-" in arrow or "-." in arrow else "solid"
        edges.append(Edge(src=src, dst=dst, label=label, kind=kind))
    return edges


def export_pptx(
    *,
    svg_path: Path,
    mmd_path: Path,
    out_path: Path,
    slide_width_in: float = 13.333,
) -> None:
    viewbox, clusters, nodes = load_nodes(svg_path)
    edges = load_edges(mmd_path)

    # Map SVG coordinates to PPTX coordinates with same aspect ratio.
    aspect = viewbox.w / viewbox.h
    slide_w = Inches(slide_width_in)
    slide_h = Inches(slide_width_in / aspect)

    prs = Presentation()
    prs.slide_width = slide_w
    prs.slide_height = slide_h

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    sx = slide_w / viewbox.w
    sy = slide_h / viewbox.h

    def map_x(x: float) -> int:
        return int((x - viewbox.x) * sx)

    def map_y(y: float) -> int:
        return int((y - viewbox.y) * sy)

    def map_w(w: float) -> int:
        return int(w * sx)

    def map_h(h: float) -> int:
        return int(h * sy)

    # Add clusters first (background)
    for c in clusters:
        left, top, w, h = map_x(c.box.x), map_y(c.box.y), map_w(c.box.w), map_h(c.box.h)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFC, 0xFC, 0xFC)
        shape.line.color.rgb = RGBColor(0x70, 0x70, 0x70)
        shape.line.width = Pt(1)
        # Cluster title
        title = slide.shapes.add_textbox(left + Pt(6), top + Pt(4), w - Pt(12), Pt(18))
        tf = title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = c.label
        run.font.bold = True
        run.font.size = Pt(11)
        run.font.name = "IBM Plex Mono"
        p.alignment = PP_ALIGN.LEFT

    # Map nodes by Mermaid key
    by_key: dict[str, Node] = {n.key: n for n in nodes}

    def node_center(n: Node) -> tuple[int, int]:
        cx = map_x(n.box.x) + map_w(n.box.w) // 2
        cy = map_y(n.box.y) + map_h(n.box.h) // 2
        return cx, cy

    # Add edges (simple straight connectors for editability)
    for e in edges:
        if e.src not in by_key or e.dst not in by_key:
            continue
        a = by_key[e.src]
        b = by_key[e.dst]
        ax, ay = node_center(a)
        bx, by = node_center(b)
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ax, ay, bx, by)
        conn.line.color.rgb = RGBColor(0x66, 0x66, 0x66)
        conn.line.width = Pt(1.25)
        if e.kind == "dotted":
            conn.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT

        if e.label:
            midx = int((ax + bx) / 2)
            midy = int((ay + by) / 2)
            lab = slide.shapes.add_textbox(midx - Pt(40), midy - Pt(10), Pt(80), Pt(20))
            lab.fill.solid()
            lab.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            lab.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            lab.line.width = Pt(0.75)
            tf = lab.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = e.label
            run.font.size = Pt(9)
            run.font.name = "IBM Plex Mono"

    # Add nodes last (foreground)
    for n in nodes:
        left, top, w, h = map_x(n.box.x), map_y(n.box.y), map_w(n.box.w), map_h(n.box.h)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, w, h)

        # Style (fill/stroke)
        if n.style.fill and n.style.fill.startswith("#"):
            shape.fill.solid()
            shape.fill.fore_color.rgb = _hex_to_rgb(n.style.fill)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)

        if n.style.stroke and n.style.stroke.startswith("#"):
            shape.line.color.rgb = _hex_to_rgb(n.style.stroke)
        else:
            shape.line.color.rgb = RGBColor(0x99, 0x99, 0x99)

        # Rough px->pt conversion for line width.
        if n.style.stroke_width_px is not None:
            shape.line.width = Pt(max(0.75, n.style.stroke_width_px * 0.75))
        else:
            shape.line.width = Pt(1)

        dash = _dash_from_svg(n.style.stroke_dasharray)
        if dash is not None:
            shape.line.dash_style = dash

        # Text
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0

        first = True
        for txt, bold in n.lines:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = txt
            run.font.name = "IBM Plex Mono"
            run.font.size = Pt(10)
            run.font.bold = bool(bold)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


def main() -> None:
    here = Path(__file__).resolve().parent
    export_pptx(
        svg_path=here / "d2ic_class_diagram_inkscape.svg",
        mmd_path=here / "d2ic_class_diagram.mmd",
        out_path=here / "d2ic_class_diagram_editable.pptx",
    )


if __name__ == "__main__":
    main()
